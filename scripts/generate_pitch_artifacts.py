"""Generate HarnessAgent pitch PDFs and PPTX decks.

The source narrative lives in docs/pitch/*.md. This script creates polished,
portable pitch artifacts under output/pitch/.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import (
    Flowable,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "pitch"

INK = "#111827"
MUTED = "#526070"
SOFT = "#EFF4FA"
PANEL = "#FFFFFF"
BORDER = "#CBD5E1"
BLUE = "#2563EB"
TEAL = "#0F766E"
AMBER = "#B45309"
GREEN = "#15803D"
RED = "#B91C1C"
DARK = "#0B1220"


def hex_color(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def reportlab_color(value: str):
    return colors.HexColor(value)


class ArchitectureDiagram(Flowable):
    """ReportLab architecture diagram for the technical PDF."""

    def __init__(self, width: float = 7.25 * inch, height: float = 3.45 * inch):
        super().__init__()
        self.width = width
        self.height = height

    def draw(self) -> None:
        c = self.canv
        c.setStrokeColor(reportlab_color(BORDER))
        c.setLineWidth(1)
        c.setFillColor(reportlab_color("#F8FAFC"))
        c.roundRect(0, 0, self.width, self.height, 10, fill=1, stroke=1)

        boxes = {
            "UI / API": (18, 166, 100, 38, BLUE),
            "Redis / RQ": (142, 166, 100, 38, AMBER),
            "Worker": (266, 166, 100, 38, TEAL),
            "AgentRunner": (390, 166, 100, 38, BLUE),
            "BaseAgent": (205, 95, 112, 38, DARK),
            "LLM Router": (28, 30, 96, 34, BLUE),
            "Memory": (150, 30, 96, 34, TEAL),
            "Tools": (272, 30, 96, 34, AMBER),
            "Guardrails": (394, 30, 96, 34, RED),
        }
        for label, (x, y, w, h, fill) in boxes.items():
            c.setFillColor(reportlab_color(fill))
            c.setStrokeColor(reportlab_color(fill))
            c.roundRect(x, y, w, h, 7, fill=1, stroke=1)
            c.setFillColor(colors.white)
            c.setFont("Helvetica-Bold", 8.5)
            c.drawCentredString(x + w / 2, y + h / 2 - 3, label)

        arrows = [
            ("UI / API", "Redis / RQ"),
            ("Redis / RQ", "Worker"),
            ("Worker", "AgentRunner"),
            ("AgentRunner", "BaseAgent"),
            ("BaseAgent", "LLM Router"),
            ("BaseAgent", "Memory"),
            ("BaseAgent", "Tools"),
            ("BaseAgent", "Guardrails"),
        ]
        for start, end in arrows:
            sx, sy, sw, sh, _ = boxes[start]
            ex, ey, ew, eh, _ = boxes[end]
            x1, y1 = sx + sw / 2, sy
            x2, y2 = ex + ew / 2, ey + eh
            if sy == ey:
                x1, y1 = sx + sw, sy + sh / 2
                x2, y2 = ex, ey + eh / 2
            c.setStrokeColor(reportlab_color("#64748B"))
            c.line(x1, y1, x2, y2)
            c.circle(x2, y2, 2, fill=1, stroke=0)

        c.setFillColor(reportlab_color(INK))
        c.setFont("Helvetica-Bold", 10)
        c.drawString(18, self.height - 24, "Production agent control plane")
        c.setFillColor(reportlab_color(MUTED))
        c.setFont("Helvetica", 8)
        c.drawString(18, self.height - 42, "Runs, handoffs, tools, guardrails, evaluation, and cost are measured in one lifecycle.")


def style_sheet():
    styles = getSampleStyleSheet()
    styles["Title"].fontName = "Helvetica-Bold"
    styles["Title"].fontSize = 24
    styles["Title"].leading = 28
    styles["Title"].textColor = reportlab_color(INK)
    styles["Heading1"].fontName = "Helvetica-Bold"
    styles["Heading1"].fontSize = 15
    styles["Heading1"].leading = 18
    styles["Heading1"].spaceBefore = 10
    styles["Heading1"].spaceAfter = 8
    styles["Heading1"].textColor = reportlab_color(INK)
    styles["Heading2"].fontName = "Helvetica-Bold"
    styles["Heading2"].fontSize = 12
    styles["Heading2"].leading = 15
    styles["Heading2"].textColor = reportlab_color(INK)
    styles["BodyText"].fontName = "Helvetica"
    styles["BodyText"].fontSize = 9.5
    styles["BodyText"].leading = 13
    styles["BodyText"].textColor = reportlab_color("#1F2937")
    return styles


def para(text: str, style_name: str = "BodyText") -> Paragraph:
    return Paragraph(text, style_sheet()[style_name])


def bullets(items: list[str]) -> list[Paragraph]:
    return [para(f"- {item}") for item in items]


def table(data: list[list[str]], widths: list[float] | None = None) -> Table:
    tbl = Table(data, colWidths=widths, hAlign="LEFT")
    tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), reportlab_color(DARK)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("LEADING", (0, 0), (-1, -1), 10),
                ("TEXTCOLOR", (0, 1), (-1, -1), reportlab_color("#1F2937")),
                ("GRID", (0, 0), (-1, -1), 0.4, reportlab_color(BORDER)),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, reportlab_color("#F8FAFC")]),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )
    return tbl


def build_pdf(path: Path, title: str, subtitle: str, sections: list[dict]) -> None:
    doc = SimpleDocTemplate(
        str(path),
        pagesize=LETTER,
        rightMargin=0.62 * inch,
        leftMargin=0.62 * inch,
        topMargin=0.58 * inch,
        bottomMargin=0.58 * inch,
        title=title,
    )
    story = [para(title, "Title"), para(subtitle), Spacer(1, 0.18 * inch)]
    for section in sections:
        if section.get("page_break"):
            story.append(PageBreak())
        story.append(para(section["heading"], "Heading1"))
        for block in section["blocks"]:
            kind = block[0]
            if kind == "p":
                story.append(para(block[1]))
                story.append(Spacer(1, 0.08 * inch))
            elif kind == "bullets":
                story.extend(bullets(block[1]))
                story.append(Spacer(1, 0.08 * inch))
            elif kind == "table":
                story.append(table(block[1], block[2] if len(block) > 2 else None))
                story.append(Spacer(1, 0.12 * inch))
            elif kind == "diagram":
                story.append(ArchitectureDiagram())
                story.append(Spacer(1, 0.12 * inch))
    doc.build(story, onFirstPage=footer, onLaterPages=footer)


def footer(c: canvas.Canvas, doc) -> None:
    c.saveState()
    c.setStrokeColor(reportlab_color(BORDER))
    c.line(doc.leftMargin, 0.42 * inch, LETTER[0] - doc.rightMargin, 0.42 * inch)
    c.setFillColor(reportlab_color(MUTED))
    c.setFont("Helvetica", 7.5)
    c.drawString(doc.leftMargin, 0.25 * inch, "HarnessAgent pitch material")
    c.drawRightString(LETTER[0] - doc.rightMargin, 0.25 * inch, f"Page {doc.page}")
    c.restoreState()


def add_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_color("#F8FAFC")
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.52))
    box.fill.solid()
    box.fill.fore_color.rgb = hex_color(DARK)
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.18), Inches(4.5), Inches(0.2))
    tf = tb.text_frame
    tf.text = "HarnessAgent"
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.78), Inches(8.4), Inches(0.72))
    title_box.text_frame.text = title
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_color(INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.56), Inches(1.45), Inches(7.8), Inches(0.45))
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(12)
        sub.text_frame.paragraphs[0].font.color.rgb = hex_color(MUTED)
    return slide


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 14, bold: bool = False, color: str = INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = hex_color(color)
    return shape


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, accent: str = BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.color.rgb = hex_color(BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_color(accent)
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.32, 0.25, 12, True, INK)
    add_text(slide, body, x + 0.18, y + 0.52, w - 0.32, h - 0.62, 9, False, MUTED)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(13)
        p.font.color.rgb = hex_color(INK)
        p.space_after = Pt(8)


def add_architecture_slide(slide):
    nodes = [
        ("UI/API", 0.75, 2.0, BLUE),
        ("Redis/RQ", 2.55, 2.0, AMBER),
        ("Worker", 4.35, 2.0, TEAL),
        ("AgentRunner", 6.15, 2.0, BLUE),
        ("BaseAgent", 4.35, 3.35, DARK),
        ("LLM Router", 0.95, 5.0, BLUE),
        ("Memory", 3.2, 5.0, TEAL),
        ("Tools", 5.45, 5.0, AMBER),
        ("Guardrails", 7.7, 5.0, RED),
        ("Eval/UI Gates", 9.95, 5.0, GREEN),
    ]
    centers = {}
    for label, x, y, color in nodes:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_color(color)
        shape.line.fill.background()
        shape.text_frame.text = label
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        centers[label] = (x + 0.675, y + 0.275)
    for a, b in [
        ("UI/API", "Redis/RQ"),
        ("Redis/RQ", "Worker"),
        ("Worker", "AgentRunner"),
        ("AgentRunner", "BaseAgent"),
        ("BaseAgent", "LLM Router"),
        ("BaseAgent", "Memory"),
        ("BaseAgent", "Tools"),
        ("BaseAgent", "Guardrails"),
        ("BaseAgent", "Eval/UI Gates"),
    ]:
        x1, y1 = centers[a]
        x2, y2 = centers[b]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = hex_color("#64748B")
        line.line.width = Pt(1.2)


def deck(path: Path, mode: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if mode == "technical":
        s = add_slide(prs, "HarnessAgent technical pitch", "A production control plane for single-agent and multi-agent AI systems.")
        add_card(s, 0.7, 2.2, 2.8, 2.2, "Core thesis", "Agents need a runtime harness: tools, memory, routing, guardrails, traces, cost, evals, and approvals.", BLUE)
        add_card(s, 3.9, 2.2, 2.8, 2.2, "Latest addition", "Single and multi-agent evals now classify failures and expose tool, guardrail, handoff, cache, token, and cost metrics.", TEAL)
        add_card(s, 7.1, 2.2, 2.8, 2.2, "Production focus", "The UI, API, and worker contracts now expose only runnable native agents while adapter factories mature.", GREEN)

        s = add_slide(prs, "Architecture: one lifecycle for run, guard, observe, evaluate", "The harness turns agent execution into an auditable control-plane flow.")
        add_architecture_slide(s)

        s = add_slide(prs, "Why agent systems fail without a harness")
        add_bullets(s, [
            "Tool calls break at schema, timeout, permission, or safety boundaries.",
            "Multi-agent handoffs lose context unless predecessor outputs are explicitly measured.",
            "Model choice and retrieved context inflate cost without clear quality gains.",
            "Prompt fixes ship without regression gates.",
            "Operators cannot see where guardrails block, pass, or need HITL.",
        ], 0.8, 1.75, 8.2, 4.8)

        s = add_slide(prs, "Evaluation pipeline: single-agent and multi-agent gates")
        add_card(s, 0.7, 1.8, 2.4, 2.0, "Single smoke", "Runs SQL/code cases through AgentRunner and reports pass rate, cost, tokens, tools, cache.", BLUE)
        add_card(s, 3.35, 1.8, 2.4, 2.0, "Multi handoff", "Runs TaskPlan DAGs through Scheduler and records sub-agent handoff metrics.", TEAL)
        add_card(s, 6.0, 1.8, 2.4, 2.0, "Prompt compare", "Compares baseline and patched prompt versions with quality and cost deltas.", AMBER)
        add_card(s, 8.65, 1.8, 2.4, 2.0, "Diagnostics", "Classifies tool, safety, budget, LLM, memory, router, planner, communication, and quality failures.", GREEN)

        s = add_slide(prs, "What operators can see")
        add_card(s, 0.7, 1.7, 2.4, 1.7, "Tools", "Calls and errors by agent.", AMBER)
        add_card(s, 3.35, 1.7, 2.4, 1.7, "Guardrails", "Blocked or HITL-routed checks.", RED)
        add_card(s, 6.0, 1.7, 2.4, 1.7, "Handoffs", "Predecessor context passed to dependent agents.", TEAL)
        add_card(s, 8.65, 1.7, 2.4, 1.7, "Cost/cache", "Token hotspots and cache reuse.", BLUE)
        add_text(s, "The dashboard now calls /evals endpoints and renders live metrics in the Latest Eval Run panel.", 0.8, 4.25, 9.8, 0.7, 17, True, INK)

        s = add_slide(prs, "Technical roadmap")
        add_bullets(s, [
            "Persist eval datasets, reports, and trends for CI/CD release gates.",
            "Complete framework adapter factories for LangGraph, CrewAI, AutoGen, and OpenClaw.",
            "Add router policy UI for cost/quality model selection.",
            "Expand enterprise controls: RBAC, audit exports, data residency, policy templates.",
        ], 0.8, 1.75, 9.8, 4.5)
    else:
        s = add_slide(prs, "HarnessAgent business pitch", "The operating layer that makes enterprise AI agents measurable, safe, and cost-controlled.")
        add_card(s, 0.7, 2.1, 2.7, 2.1, "Buyer pain", "Agent demos are easy. Agent operations are hard: trust, cost, risk, and regression control.", RED)
        add_card(s, 3.8, 2.1, 2.7, 2.1, "Product wedge", "Start with SQL and code agents, then expand into multi-agent workflows.", BLUE)
        add_card(s, 6.9, 2.1, 2.7, 2.1, "Business value", "Ship agents faster while lowering model spend and incident risk.", GREEN)

        s = add_slide(prs, "Market story: governed agents beat unmanaged demos")
        add_bullets(s, [
            "CTOs need production readiness, not one-off automations.",
            "Risk teams need proof that guardrails and approvals are working.",
            "FinOps needs spend attribution by agent, tenant, task, and model.",
            "Product teams need prompt improvements backed by eval evidence.",
        ], 0.8, 1.75, 8.8, 4.6)

        s = add_slide(prs, "Business architecture: build, run, guard, evaluate, improve")
        add_architecture_slide(s)
        add_text(s, "One control plane connects engineering velocity to risk controls and cost discipline.", 0.9, 6.35, 9.4, 0.45, 15, True, INK)

        s = add_slide(prs, "Stakeholder value")
        add_card(s, 0.7, 1.6, 2.3, 1.6, "CTO", "Production readiness and governance.", BLUE)
        add_card(s, 3.25, 1.6, 2.3, 1.6, "Engineering", "Shared runtime instead of duplicated platform plumbing.", TEAL)
        add_card(s, 5.8, 1.6, 2.3, 1.6, "Risk", "Guardrails, HITL, audit, blocked actions.", RED)
        add_card(s, 8.35, 1.6, 2.3, 1.6, "Finance", "Token, cost, and cache attribution.", GREEN)

        s = add_slide(prs, "Pilot plan")
        add_card(s, 0.75, 1.6, 2.0, 3.4, "Weeks 1-2", "Connect one SQL/code agent. Add tool contracts and guardrails.", BLUE)
        add_card(s, 3.1, 1.6, 2.0, 3.4, "Weeks 3-4", "Build smoke/regression evals. Deploy dashboard to operators.", TEAL)
        add_card(s, 5.45, 1.6, 2.0, 3.4, "Weeks 5-6", "Optimize prompt, router, cache. Executive readout with metrics.", GREEN)
        add_text(s, "Pilot evidence: pass rate, failure-stage mix, average cost, guardrail hits, tool errors, handoff coverage.", 8.0, 2.0, 3.7, 2.4, 15, True, INK)

        s = add_slide(prs, "Investment ask")
        add_bullets(s, [
            "Fund eval history and CI release gates.",
            "Productize enterprise access controls and audit exports.",
            "Complete adapter factories for broader agent adoption.",
            "Add FinOps dashboards for routing and spend optimization.",
        ], 0.8, 1.75, 9.0, 4.5)

    prs.save(path)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    build_pdf(
        OUT / "HarnessAgent_Technical_Pitch.pdf",
        "HarnessAgent Technical Pitch",
        "Production control plane for single-agent and multi-agent AI systems.",
        [
            {
                "heading": "Executive summary",
                "blocks": [
                    ("p", "HarnessAgent turns agent prototypes into governed systems with routing, memory, tools, guardrails, HITL, observability, cost controls, and evaluation gates."),
                    ("p", "The latest implementation adds single-agent and multi-agent eval diagnostics with tool, guardrail, handoff, cache, token, and cost visibility."),
                ],
            },
            {
                "heading": "Architecture diagram",
                "blocks": [("diagram",)],
            },
            {
                "heading": "What the harness measures",
                "blocks": [
                    (
                        "table",
                        [
                            ["Metric", "Operational value"],
                            ["Failure stage", "Separates tool, safety, budget, LLM, memory, router, planner, communication, and quality failures."],
                            ["Tool calls/errors", "Shows schema, timeout, missing tool, or unsafe execution problems."],
                            ["Guardrail hits", "Shows policy blocks and HITL routing."],
                            ["Handoffs", "Makes multi-agent communication measurable."],
                            ["Cost/cache", "Finds token hotspots and context reuse opportunities."],
                        ],
                        [1.55 * inch, 5.45 * inch],
                    )
                ],
            },
            {
                "heading": "Production readiness",
                "page_break": True,
                "blocks": [
                    ("bullets", [
                        "Redis/RQ-backed execution and persisted run records.",
                        "Runnable-agent contract aligned across UI, API, and worker.",
                        "Safety checks around model, tool, output, budget, and HITL boundaries.",
                        "UI-triggered /evals endpoints for smoke, multi-agent regression, and prompt comparison.",
                    ])
                ],
            },
            {
                "heading": "Technical roadmap",
                "blocks": [
                    ("bullets", [
                        "Persist eval datasets and reports for CI/CD gates.",
                        "Complete framework adapter factories.",
                        "Add router policy UI and enterprise audit controls.",
                    ])
                ],
            },
        ],
    )
    build_pdf(
        OUT / "HarnessAgent_Business_Pitch.pdf",
        "HarnessAgent Business Pitch",
        "The operating layer for governed, measurable, cost-controlled AI agents.",
        [
            {
                "heading": "Positioning",
                "blocks": [
                    ("p", "HarnessAgent is not another chatbot builder. It is the harness around agents: evaluation, guardrails, tools, memory, routing, observability, cost controls, and human approvals."),
                    ("p", "The business value is simple: deploy useful agents faster while reducing trust, compliance, and model-spend risk."),
                ],
            },
            {
                "heading": "Business architecture",
                "blocks": [("diagram",)],
            },
            {
                "heading": "Stakeholder value",
                "blocks": [
                    (
                        "table",
                        [
                            ["Stakeholder", "Need", "HarnessAgent answer"],
                            ["CTO / AI leader", "Production readiness", "One control plane for runs, tools, evals, and policies."],
                            ["Engineering", "Faster platform delivery", "Shared runtime for routing, memory, tools, telemetry, and safety."],
                            ["Risk", "Proof of controls", "Guardrail hits, HITL queues, audit logs, blocked actions."],
                            ["Finance", "Cost attribution", "Token, cost, cache, and per-agent spend visibility."],
                        ],
                        [1.35 * inch, 2.0 * inch, 3.75 * inch],
                    )
                ],
            },
            {
                "heading": "Pilot plan",
                "page_break": True,
                "blocks": [
                    (
                        "table",
                        [
                            ["Week", "Milestone", "Evidence"],
                            ["1-2", "Connect SQL/code agent and guardrails", "First governed run with trace, cost, tool metrics."],
                            ["3-4", "Build eval suite and deploy dashboard", "Pass rate, failure stage, token baseline."],
                            ["5-6", "Optimize prompts, router, cache", "Measured cost reduction or pass-rate lift."],
                        ],
                        [0.75 * inch, 2.55 * inch, 3.8 * inch],
                    )
                ],
            },
            {
                "heading": "Investment ask",
                "blocks": [
                    ("bullets", [
                        "Fund eval history and CI/CD release gates.",
                        "Productize enterprise access controls and audit exports.",
                        "Complete framework adapter factories.",
                        "Add FinOps dashboards for routing and spend optimization.",
                    ])
                ],
            },
        ],
    )
    deck(OUT / "HarnessAgent_Technical_Pitch.pptx", "technical")
    deck(OUT / "HarnessAgent_Business_Pitch.pptx", "business")
    print(f"Generated pitch artifacts in {OUT}")


if __name__ == "__main__":
    main()
