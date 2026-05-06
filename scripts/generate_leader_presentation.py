"""Generate a Gartner-style leadership presentation for HarnessAgent."""

from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "pitch"
DECK_PATH = OUT / "HarnessAgent_Leadership_Gartner_Style_Presentation.pptx"
ARCH_IMAGE = OUT / "HarnessAgent_Leadership_Architecture.png"

INK = "#111827"
DARK = "#08111F"
MUTED = "#526070"
GRID = "#D7DEE8"
PANEL = "#FFFFFF"
BG = "#F6F8FB"
BLUE = "#2463EB"
CYAN = "#0891B2"
TEAL = "#0F766E"
AMBER = "#B45309"
GREEN = "#15803D"
RED = "#B91C1C"
VIOLET = "#6D28D9"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def pil_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_wrapped(draw: ImageDraw.ImageDraw, text: str, xy: tuple[int, int], width: int, fill: str, size: int, bold: bool = False) -> None:
    x, y = xy
    line_height = int(size * 1.25)
    chars = max(16, width // max(7, int(size * 0.55)))
    for line in wrap(text, chars):
        draw.text((x, y), line, fill=pil_rgb(fill), font=font(size, bold))
        y += line_height


def rounded_box(draw: ImageDraw.ImageDraw, rect: tuple[int, int, int, int], fill: str, outline: str = GRID, radius: int = 18) -> None:
    draw.rounded_rectangle(rect, radius=radius, fill=pil_rgb(fill), outline=pil_rgb(outline), width=2)


def arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], fill: str = MUTED) -> None:
    draw.line([start, end], fill=pil_rgb(fill), width=3)
    x1, y1 = start
    x2, y2 = end
    if abs(x2 - x1) >= abs(y2 - y1):
        pts = [(x2, y2), (x2 - 12 if x2 > x1 else x2 + 12, y2 - 7), (x2 - 12 if x2 > x1 else x2 + 12, y2 + 7)]
    else:
        pts = [(x2, y2), (x2 - 7, y2 - 12 if y2 > y1 else y2 + 12), (x2 + 7, y2 - 12 if y2 > y1 else y2 + 12)]
    draw.polygon(pts, fill=pil_rgb(fill))


def create_architecture_image() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1600, 900), pil_rgb(BG))
    d = ImageDraw.Draw(img)
    d.text((55, 42), "HarnessAgent Reference Architecture", fill=pil_rgb(INK), font=font(42, True))
    draw_wrapped(
        d,
        "A governed control plane that wraps agents with tools, memory, guardrails, evals, cost controls, and observability.",
        (58, 98),
        1180,
        MUTED,
        22,
    )

    zones = [
        ("Experience Layer", (55, 170, 360, 720), BLUE, ["Leader dashboard", "Operator console", "API clients"]),
        ("Control Plane", (425, 170, 820, 720), DARK, ["FastAPI", "Run records", "RQ queues", "Eval API"]),
        ("Agent Runtime", (885, 170, 1238, 720), TEAL, ["AgentRunner", "Scheduler", "BaseAgent", "HITL"]),
        ("Enterprise Services", (1250, 170, 1545, 720), AMBER, ["LLM providers", "Databases", "MCP servers", "Repos"]),
    ]
    for title, rect, color, items in zones:
        rounded_box(d, rect, "#FFFFFF")
        x1, y1, x2, _ = rect
        d.rounded_rectangle((x1, y1, x2, y1 + 54), radius=18, fill=pil_rgb(color))
        d.text((x1 + 22, y1 + 15), title, fill=(255, 255, 255), font=font(23, True))
        for idx, item in enumerate(items):
            y = y1 + 92 + idx * 86
            rounded_box(d, (x1 + 24, y, x2 - 24, y + 54), "#F8FAFC", "#E2E8F0", 12)
            d.text((x1 + 45, y + 16), item, fill=pil_rgb(INK), font=font(21, True))

    for start, end in [((360, 445), (425, 445)), ((820, 445), (885, 445)), ((1238, 445), (1250, 445))]:
        arrow(d, start, end)

    metrics = [
        ("Evaluate", "pass rate, failure stage", GREEN),
        ("Govern", "guardrails, HITL, audit", RED),
        ("Optimize", "tokens, cost, cache", VIOLET),
    ]
    for idx, (title, body, color) in enumerate(metrics):
        x = 480 + idx * 245
        rounded_box(d, (x, 752, x + 215, 835), "#FFFFFF", color, 16)
        d.text((x + 18, 770), title, fill=pil_rgb(color), font=font(22, True))
        d.text((x + 18, 800), body, fill=pil_rgb(MUTED), font=font(16))

    img.save(ARCH_IMAGE)


def slide(prs: Presentation, title: str, subtitle: str = "", section: str = "HarnessAgent"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BG)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.43))
    top.fill.solid()
    top.fill.fore_color.rgb = rgb(DARK)
    top.line.fill.background()
    add_text(s, section.upper(), 0.45, 0.13, 4.4, 0.2, 8, True, "#FFFFFF")
    add_text(s, title, 0.55, 0.72, 9.9, 0.65, 26, True, INK)
    if subtitle:
        add_text(s, subtitle, 0.58, 1.34, 9.6, 0.38, 11, False, MUTED)
    add_text(s, "Leadership briefing", 10.65, 0.13, 2.1, 0.2, 8, False, "#DCE6F3")
    return s


def add_text(s, text: str, x: float, y: float, w: float, h: float, size: int = 12, bold: bool = False, color: str = INK, align=PP_ALIGN.LEFT):
    shape = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return shape


def add_label(s, text: str, x: float, y: float, w: float, h: float, size: int = 12, bold: bool = False, color: str = INK, align=PP_ALIGN.LEFT):
    shape = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return shape


def card(s, x: float, y: float, w: float, h: float, title: str, body: str, accent: str = BLUE):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(PANEL)
    shape.line.color.rgb = rgb(GRID)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    add_label(s, title, x + 0.2, y + 0.18, w - 0.35, 0.38, 11, True, INK)
    add_label(s, body, x + 0.2, y + 0.68, w - 0.35, h - 0.82, 8.5, False, MUTED)


def metric(s, x: float, y: float, value: str, label: str, note: str, color: str = BLUE):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.35), Inches(1.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(PANEL)
    shape.line.color.rgb = rgb(GRID)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.34))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color)
    bar.line.fill.background()
    add_label(s, value, x + 0.24, y + 0.15, 1.65, 0.38, 22, True, color)
    add_label(s, label, x + 0.24, y + 0.6, 1.9, 0.24, 10, True, INK)
    add_label(s, note, x + 0.24, y + 0.9, 1.88, 0.32, 8.2, False, MUTED)


def bullets(s, items: list[str], x: float, y: float, w: float, h: float, size: int = 13):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(INK)
        p.space_after = Pt(7)


def connector(s, x1: float, y1: float, x2: float, y2: float):
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb("#64748B")
    line.line.width = Pt(1.3)


def build_deck() -> None:
    create_architecture_image()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = slide(
        prs,
        "HarnessAgent: governed operating layer for enterprise AI agents",
        "A leadership briefing on business need, architecture, features, workflow, and API surface.",
        "Executive summary",
    )
    metric(s, 0.75, 2.05, "01", "Control plane", "Standardize how agents run, fail, improve, and get approved.", BLUE)
    metric(s, 3.35, 2.05, "02", "Risk posture", "Make guardrails, HITL, tool actions, and audit evidence visible.", RED)
    metric(s, 5.95, 2.05, "03", "Cost discipline", "Attribute tokens, cost, cache, and expensive agents.", GREEN)
    add_text(s, "Leadership decision", 0.78, 4.15, 3.2, 0.3, 13, True, DARK)
    add_text(
        s,
        "Move from scattered agent demos to a governed platform capability: build once, operate many, promote only when eval gates pass.",
        0.78,
        4.55,
        9.2,
        0.75,
        20,
        True,
        INK,
    )

    s = slide(prs, "The business problem: agent demos do not equal agent operations", "The enterprise risk is not model access. It is unmanaged execution.", "Business problem")
    card(s, 0.7, 1.85, 2.7, 1.9, "Trust gap", "Leaders cannot see which agents are reliable enough for real workflows.", RED)
    card(s, 3.65, 1.85, 2.7, 1.9, "Cost gap", "Token spend grows without attribution by agent, tenant, model, or task.", AMBER)
    card(s, 6.6, 1.85, 2.7, 1.9, "Control gap", "Tool calls, guardrails, and human approvals are fragmented.", BLUE)
    card(s, 9.55, 1.85, 2.7, 1.9, "Quality gap", "Prompt changes ship without regression or multi-agent handoff tests.", VIOLET)
    add_text(s, "Implication", 0.8, 4.45, 2.0, 0.3, 13, True, DARK)
    add_text(s, "Without a harness, each team rebuilds the same runtime controls and still cannot prove agent quality, safety, or cost efficiency.", 0.8, 4.85, 9.5, 0.7, 19, True, INK)

    s = slide(prs, "Reference architecture: a governed AI agent control plane", "Architecture diagram using a generated image asset embedded in the deck.", "Architecture")
    s.shapes.add_picture(str(ARCH_IMAGE), Inches(0.65), Inches(1.55), width=Inches(12.0))

    s = slide(prs, "Capability map: what leaders get from the platform", "Capabilities map to executive controls: quality, risk, cost, and velocity.", "Features")
    caps = [
        ("Run lifecycle", "Create, queue, execute, stream, cancel, and inspect runs.", BLUE),
        ("Evaluation", "Single-agent smoke, multi-agent handoff, prompt comparison.", GREEN),
        ("Guardrails", "Safety policies, destructive-tool approval, HITL, audit context.", RED),
        ("Tooling", "Schema-validated SQL, code, file, and MCP tools.", AMBER),
        ("Model routing", "Provider fallback, context-window fit, local/cloud options.", CYAN),
        ("Observability", "Events, failure stages, token/cost metrics, MLflow/OTel.", VIOLET),
    ]
    for i, (title, body, color) in enumerate(caps):
        x = 0.72 + (i % 3) * 4.05
        y = 1.7 + (i // 3) * 2.0
        card(s, x, y, 3.55, 1.55, title, body, color)

    s = slide(prs, "How it works: controlled execution loop", "Each run moves through an observable lifecycle with release-gate feedback.", "Operating model")
    steps = [
        ("1. Request", "Operator or API creates run"),
        ("2. Route", "Queue selects worker and agent"),
        ("3. Execute", "Agent calls model, memory, tools"),
        ("4. Guard", "Policies and HITL control risk"),
        ("5. Evaluate", "Diagnostics classify failures"),
        ("6. Improve", "Prompt/router/cache changes gated"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.75 + i * 2.0
        card(s, x, 2.2, 1.58, 1.35, title, body, [BLUE, CYAN, TEAL, RED, GREEN, VIOLET][i])
        if i < len(steps) - 1:
            connector(s, x + 1.58, 2.88, x + 1.96, 2.88)
    add_text(s, "Feedback loop", 0.8, 4.5, 1.7, 0.3, 13, True, DARK)
    add_text(s, "Every failed run becomes structured learning: failure stage, cost driver, guardrail action, tool error, handoff gap, and optimization hint.", 0.8, 4.92, 9.8, 0.55, 18, True, INK)

    s = slide(prs, "API surface: endpoints leaders should know exist", "The API exposes the core governance lifecycle and can be wrapped by portals, chatops, or CI gates.", "API and endpoints")
    endpoint_rows = [
        ("POST /runs", "Create and enqueue a governed agent run."),
        ("GET /runs", "List tenant-scoped run history."),
        ("GET /runs/{run_id}", "Inspect status, result, metadata, and HITL state."),
        ("GET /runs/{run_id}/stream", "Stream step events and token deltas."),
        ("DELETE /runs/{run_id}", "Cancel a pending or running execution."),
        ("POST /evals/smoke/run", "Run built-in single-agent smoke gates."),
        ("POST /evals/multi/run", "Run multi-agent handoff regression gates."),
        ("POST /evals/compare", "Compare baseline and patched prompt versions."),
        ("GET /health", "Check service and dependency readiness."),
        ("GET /memory/*", "Manage and inspect memory when enabled."),
    ]
    for i, (path, desc) in enumerate(endpoint_rows):
        col = i // 5
        row = i % 5
        x = 0.78 + col * 5.9
        y = 1.62 + row * 0.86
        card(s, x, y, 5.2, 0.62, path, desc, BLUE if col == 0 else GREEN)

    s = slide(prs, "Decision framework: where HarnessAgent creates leverage", "Use this as an executive scorecard for pilot success.", "Value case")
    metric(s, 0.8, 1.75, "Quality", "Release gates", "Pass-rate and failure-stage trends.", GREEN)
    metric(s, 3.4, 1.75, "Risk", "Policy evidence", "Guardrails and HITL queues.", RED)
    metric(s, 6.0, 1.75, "Cost", "FinOps visibility", "Tokens, cache, and model routing.", AMBER)
    metric(s, 8.6, 1.75, "Speed", "Shared runtime", "Reusable safety, tools, and tracing.", BLUE)
    bullets(
        s,
        [
            "Pilot target: one SQL agent, one code agent, one multi-agent handoff workflow.",
            "Promotion rule: pass smoke and handoff gates with no critical guardrail regressions.",
            "Leadership readout: quality, cost, risk, adoption, and roadmap decisions in one dashboard.",
        ],
        0.9,
        4.1,
        9.5,
        1.6,
        14,
    )

    s = slide(prs, "Recommended next steps", "A practical adoption path for leaders and platform teams.", "Roadmap")
    card(s, 0.8, 1.75, 2.25, 3.25, "0-30 days", "Run a pilot with SQL/code agents, built-in evals, and operator dashboard.", BLUE)
    card(s, 3.35, 1.75, 2.25, 3.25, "30-60 days", "Persist eval history, connect CI gates, tune router and cache policies.", TEAL)
    card(s, 5.9, 1.75, 2.25, 3.25, "60-90 days", "Add adapter factories, RBAC, audit exports, and enterprise guardrail packs.", VIOLET)
    card(s, 8.45, 1.75, 2.8, 3.25, "Leader ask", "Treat agent governance as platform infrastructure, not project-by-project glue code.", GREEN)

    prs.save(DECK_PATH)


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    build_deck()
    print(DECK_PATH)
